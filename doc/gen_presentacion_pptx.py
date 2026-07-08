#!/usr/bin/env python3
"""
Genera doc/PRESENTACION-SLIDES.pptx a partir de doc/PRESENTACION-SLIDES.md

Formato esperado del markdown:
  ## Slide N — Título
  **Subtítulo en negrita** (opcional)
  - viñeta
  - viñeta
  | tabla | opcional |
  > Nota: texto de la nota del ponente (puede ocupar varias líneas con '>')

Uso:
  python doc/gen_presentacion_pptx.py
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
SRC_MD = HERE / "PRESENTACION-SLIDES.md"
OUT_PPTX = HERE / "PRESENTACION-SLIDES.pptx"

# --- Paleta ---
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def parse_slides(md_text):
    """Divide el markdown en bloques de slide y extrae título, subtítulo, cuerpo y nota."""
    blocks = re.split(r"^## Slide \d+ — ", md_text, flags=re.MULTILINE)[1:]
    slides = []
    for block in blocks:
        lines = block.splitlines()
        title = lines[0].strip()
        body_lines = lines[1:]

        # Cortar en el separador '---' si existe al final
        text = "\n".join(body_lines)
        text = text.split("\n---", 1)[0]

        # Extraer nota del ponente (líneas que empiezan por '>')
        note_lines = [l[1:].strip() for l in text.splitlines() if l.strip().startswith(">")]
        note = " ".join(n[len("Nota:"):].strip() if n.startswith("Nota:") else n for n in note_lines).strip()

        # Quitar líneas de nota del cuerpo
        content_lines = [l for l in text.splitlines() if not l.strip().startswith(">")]

        subtitle = None
        bullets = []
        table_rows = []
        in_table = False
        for l in content_lines:
            s = l.strip()
            if not s:
                continue
            if s.startswith("|"):
                in_table = True
                cells = [c.strip() for c in s.strip("|").split("|")]
                if set("".join(cells).replace("-", "").strip()) == set():
                    continue  # separator row  ---|---
                table_rows.append(cells)
                continue
            if s.startswith("- "):
                bullets.append(re.sub(r"^\d+\.\s*", "", s[2:]).strip())
                continue
            m = re.match(r"^\d+\.\s+(.*)", s)
            if m:
                bullets.append(m.group(1).strip())
                continue
            if s.startswith("**") and subtitle is None and not bullets:
                subtitle = strip_md_bold(s).strip()
                continue
            # texto suelto (párrafo) -> tratar como bullet simple
            bullets.append(s)

        slides.append({
            "title": title,
            "subtitle": subtitle,
            "bullets": bullets,
            "table": table_rows,
            "note": note,
        })
    return slides


def strip_md_bold(s):
    return re.sub(r"\*\*(.*?)\*\*", r"\1", s)


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_accent_bar(slide, top=Inches(1.25), height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(1.4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def set_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def add_title(slide, title, top=Inches(0.45), size=32, color=NAVY, width=Inches(12.1)):
    box = slide.shapes.add_textbox(Inches(0.6), top, width, Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return box


def add_subtitle(slide, subtitle, top=Inches(1.45), size=20, color=DARK_TEXT):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = strip_md_bold(subtitle)
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets, top, height=Inches(4.3), size=18, color=DARK_TEXT):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = "•  " + strip_md_bold(b)
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_table(slide, rows, top, size=15):
    if not rows:
        return
    n_rows = len(rows)
    n_cols = len(rows[0])
    width = Inches(12.1)
    height = Inches(min(0.55 * n_rows, 4.0))
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, Inches(0.6), top, width, height)
    table = graphic_frame.table
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = strip_md_bold(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(size)
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE


def build_title_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    # Marca de acento
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.55), Inches(0.18), Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_data["subtitle"] or slide_data["title"]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WHITE

    sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.6))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = slide_data["bullets"][0] if slide_data["bullets"] else ""
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0xC7, 0xD3, 0xE8)

    meta_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.55), Inches(11.5), Inches(1.2))
    tf3 = meta_box.text_frame
    tf3.word_wrap = True
    for i, b in enumerate(slide_data["bullets"][1:]):
        p3 = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        run3 = p3.add_run()
        run3.text = b
        run3.font.size = Pt(16)
        run3.font.color.rgb = RGBColor(0xC7, 0xD3, 0xE8)

    set_notes(slide, slide_data["note"])
    return slide


def build_closing_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = strip_md_bold(slide_data["subtitle"] or slide_data["title"])
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(3.35), Inches(1.4), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = slide.shapes.add_textbox(Inches(0.9), Inches(3.65), Inches(11.5), Inches(2.8))
    tf2 = box.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(slide_data["bullets"]):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.space_after = Pt(10)
        run2 = p2.add_run()
        run2.text = strip_md_bold(b)
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(0xC7, 0xD3, 0xE8)

    set_notes(slide, slide_data["note"])
    return slide


def build_content_slide(prs, slide_data, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    is_highlight = slide_data["title"].strip().startswith("⭐")

    title_color = ACCENT if is_highlight else NAVY
    add_title(slide, slide_data["title"], color=title_color)
    add_accent_bar(slide)

    # número de slide
    num_box = slide.shapes.add_textbox(Inches(12.55), Inches(7.05), Inches(0.6), Inches(0.35))
    p = num_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(index)
    run.font.size = Pt(12)
    run.font.color.rgb = MUTED

    cursor_top = Inches(1.55)
    if slide_data["subtitle"]:
        add_subtitle(slide, slide_data["subtitle"], top=cursor_top)
        cursor_top = Inches(2.35)

    if slide_data["table"]:
        add_table(slide, slide_data["table"], top=cursor_top)
        cursor_top = Emu(cursor_top + Inches(0.55 * len(slide_data["table"])) + Inches(0.35))

    if slide_data["bullets"]:
        add_bullets(slide, slide_data["bullets"], top=cursor_top)

    set_notes(slide, slide_data["note"])
    return slide


def main():
    md_text = SRC_MD.read_text(encoding="utf-8")
    slides = parse_slides(md_text)
    if not slides:
        print("No se encontraron slides en el markdown.", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, s in enumerate(slides, start=1):
        if i == 1:
            build_title_slide(prs, s)
        elif i == len(slides):
            build_closing_slide(prs, s)
        else:
            build_content_slide(prs, s, i)

    prs.save(OUT_PPTX)
    print(f"Generado: {OUT_PPTX} ({len(slides)} slides)")


if __name__ == "__main__":
    main()
